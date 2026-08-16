import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_officer_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Official Law Enforcement Theme Colors
    C_NAVY_BG      = RGBColor(11, 22, 44)       # Deep Navy Blue (#0B162C)
    C_CARD_BG      = RGBColor(19, 36, 68)       # Card Surface (#132444)
    C_CARD_BORDER  = RGBColor(41, 74, 110)      # Border (#294A6E)
    C_GOLD         = RGBColor(229, 185, 83)     # Official Police Gold (#E5B953)
    C_GREEN        = RGBColor(16, 149, 93)      # Pakistan Emerald Green (#10955D)
    C_CYAN         = RGBColor(56, 189, 248)     # Sky Cyan Accent (#38BDF8)
    C_WHITE        = RGBColor(255, 255, 255)
    C_SLATE        = RGBColor(160, 174, 192)

    logo_path = os.path.abspath("f:/NCCIA/public/images/images.jpg")

    def add_slide_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_NAVY_BG
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, subtitle_text=""):
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(10.5), Inches(1.15))
        tf = header_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Georgia"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = C_GOLD
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = "Segoe UI"
            p2.font.size = Pt(12.5)
            p2.font.bold = True
            p2.font.color.rgb = C_CYAN
            p2.space_before = Pt(4)
            
        if os.path.isfile(logo_path):
            slide.shapes.add_picture(logo_path, Inches(12.0), Inches(0.3), height=Inches(1.05))

    def add_card(slide, left, top, width, height, title="", title_color=C_GOLD, badge=""):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD_BG
        card.line.color.rgb = C_CARD_BORDER
        card.line.width = Pt(1.5)
        
        if title:
            tb = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.18), Inches(width - 0.5), Inches(0.45))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = "Segoe UI"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = title_color
            
            if badge:
                p2 = tf.add_paragraph()
                p2.text = badge
                p2.font.name = "Segoe UI"
                p2.font.size = Pt(10)
                p2.font.bold = True
                p2.font.color.rgb = C_GREEN
        return card

    # ==========================================
    # SLIDE 1: WORLD-CLASS PRESIDENTIAL HERO TITLE SLIDE
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s1)

    # Decorative Outer Frame
    outer_card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.5), Inches(12.133), Inches(6.5))
    outer_card.fill.solid()
    outer_card.fill.fore_color.rgb = RGBColor(14, 28, 54)
    outer_card.line.color.rgb = RGBColor(41, 74, 110)
    outer_card.line.width = Pt(2.0)

    # Top Official Ribbon / Government Pill
    pill = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(0.85), Inches(11.333), Inches(0.55))
    pill.fill.solid()
    pill.fill.fore_color.rgb = RGBColor(10, 20, 38)
    pill.line.color.rgb = C_GOLD
    pill.line.width = Pt(1.2)
    ptf = pill.text_frame
    ptf.word_wrap = True
    ptf.margin_top = Inches(0.08)
    pp = ptf.paragraphs[0]
    pp.alignment = PP_ALIGN.CENTER
    pp.text = "ISLAMIC REPUBLIC OF PAKISTAN  •  MINISTRY OF INTERIOR & NARCOTICS CONTROL"
    pp.font.name = "Segoe UI"
    pp.font.size = Pt(12)
    pp.font.bold = True
    pp.font.color.rgb = C_GOLD

    # Left Hero Box: Official Emblem Showcase
    emblem_card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.65), Inches(4.2), Inches(4.9))
    emblem_card.fill.solid()
    emblem_card.fill.fore_color.rgb = RGBColor(18, 34, 64)
    emblem_card.line.color.rgb = RGBColor(56, 189, 248)
    emblem_card.line.width = Pt(1.5)

    # Gold Circular Glow Ring
    ring = s1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.7), Inches(2.0), Inches(2.8), Inches(2.8))
    ring.fill.solid()
    ring.fill.fore_color.rgb = C_GOLD
    ring.line.fill.background()

    # NCCIA Official Emblem
    if os.path.isfile(logo_path):
        s1.shapes.add_picture(logo_path, Inches(1.75), Inches(2.05), width=Inches(2.7), height=Inches(2.7))

    # Badge Sub-banner
    bb_box = s1.shapes.add_textbox(Inches(1.15), Inches(5.05), Inches(3.9), Inches(1.3))
    bbtf = bb_box.text_frame
    bbtf.word_wrap = True
    bb_p1 = bbtf.paragraphs[0]
    bb_p1.alignment = PP_ALIGN.CENTER
    bb_p1.text = "OFFICIAL EMBLEM"
    bb_p1.font.name = "Segoe UI"
    bb_p1.font.size = Pt(12)
    bb_p1.font.bold = True
    bb_p1.font.color.rgb = C_GOLD

    bb_p2 = bbtf.add_paragraph()
    bb_p2.alignment = PP_ALIGN.CENTER
    bb_p2.text = "Combating Cyber Crimes\nFederal Investigation Authority"
    bb_p2.font.name = "Segoe UI"
    bb_p2.font.size = Pt(10)
    bb_p2.font.color.rgb = C_SLATE
    bb_p2.space_before = Pt(3)

    # Right Content Showcase Box
    right_card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.45), Inches(1.65), Inches(6.88), Inches(4.9))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = RGBColor(18, 34, 64)
    right_card.line.color.rgb = RGBColor(41, 74, 110)
    right_card.line.width = Pt(1.5)

    # Right Content Text
    rc_box = s1.shapes.add_textbox(Inches(5.75), Inches(1.85), Inches(6.3), Inches(3.2))
    rctf = rc_box.text_frame
    rctf.word_wrap = True

    rp1 = rctf.paragraphs[0]
    rp1.text = "NATIONAL CYBER CRIME\nINVESTIGATION AGENCY"
    rp1.font.name = "Georgia"
    rp1.font.size = Pt(26)
    rp1.font.bold = True
    rp1.font.color.rgb = C_GOLD

    rp2 = rctf.add_paragraph()
    rp2.text = "قومی سائبر کرائم تفتیش ایجنسی — حکومت پاکستان"
    rp2.font.name = "Segoe UI"
    rp2.font.size = Pt(15)
    rp2.font.bold = True
    rp2.font.color.rgb = C_WHITE
    rp2.space_before = Pt(6)

    # Divider line
    div_line = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.75), Inches(3.55), Inches(6.25), Inches(0.04))
    div_line.fill.solid()
    div_line.fill.fore_color.rgb = C_GOLD
    div_line.line.fill.background()

    # System Name & Subtitle
    sys_box = s1.shapes.add_textbox(Inches(5.75), Inches(3.68), Inches(6.3), Inches(1.4))
    stf = sys_box.text_frame
    stf.word_wrap = True
    sp1 = stf.paragraphs[0]
    sp1.text = "DIGITAL CASE MANAGEMENT SYSTEM (CMS)"
    sp1.font.name = "Segoe UI"
    sp1.font.size = Pt(18)
    sp1.font.bold = True
    sp1.font.color.rgb = C_CYAN

    sp2 = stf.add_paragraph()
    sp2.text = "Standard Operating Procedures (SOPs) & Officer Operational Training Manual"
    sp2.font.name = "Segoe UI"
    sp2.font.size = Pt(12)
    sp2.font.bold = True
    sp2.font.color.rgb = C_WHITE
    sp2.space_before = Pt(4)

    # 3 Bottom Feature Pills on the Right
    pill_data = [
        ("📝 Front Desk & 80mm Slips", C_CYAN, 5.75),
        ("⚖️ VO & IO Case Enquiries", C_GOLD, 7.85),
        ("🔬 Digital Forensic Lab", C_GREEN, 9.95)
    ]
    for p_text, p_col, p_left in pill_data:
        spill = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(p_left), Inches(5.65), Inches(2.0), Inches(0.65))
        spill.fill.solid()
        spill.fill.fore_color.rgb = RGBColor(10, 22, 44)
        spill.line.color.rgb = p_col
        spill.line.width = Pt(1.2)
        sptf = spill.text_frame
        sptf.word_wrap = True
        sptf.margin_top = Inches(0.08)
        spp = sptf.paragraphs[0]
        spp.alignment = PP_ALIGN.CENTER
        spp.text = p_text
        spp.font.name = "Segoe UI"
        spp.font.size = Pt(9.5)
        spp.font.bold = True
        spp.font.color.rgb = C_WHITE

    # SLIDE 2: Purpose & Vision
    s2 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s2)
    add_header(s2, "Purpose & Vision of the CMS Portal", "Transforming Traditional Paper Policing into Transparent Digital Workflows")
    p_cards = [
        ("📁 100% Digital Case Record", "Eliminates physical file loss or delay. Every complaint, evidence, notice, and diary entry is permanently secured in the central digital database.", C_CYAN),
        ("⏱️ Fast Citizen Registration", "Front desk registration takes less than 60 seconds with instant 80mm thermal slip generation and automated WhatsApp alert dispatch.", C_GOLD),
        ("⚖️ Officer Accountability", "Complete timeline tracking from initial filing to scrutiny, verification, enquiry, and court trial. Zero hidden pendency.", C_GREEN),
        ("🔍 Instant Case Search", "Find any past or present case instantly by Complaint Number, Complainant CNIC, Phone Number, Accused Bank Account, or Offence Category.", C_CYAN)
    ]
    for idx, (title, desc, col) in enumerate(p_cards):
        row, col_idx = idx // 2, idx % 2
        c_left, c_top = 0.8 + col_idx * 5.98, 1.65 + row * 2.7
        add_card(s2, c_left, c_top, 5.75, 2.45, title, col)
        sb = s2.shapes.add_textbox(Inches(c_left + 0.25), Inches(c_top + 0.65), Inches(5.25), Inches(1.65))
        stf = sb.text_frame
        stf.word_wrap = True
        p = stf.paragraphs[0]
        p.text = desc
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.color.rgb = C_WHITE

    # SLIDE 3: Officer Roles Matrix
    s3 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s3)
    add_header(s3, "Officer Roles & Responsibilities Matrix", "Every Officer Has a Dedicated Dashboard Tailored to Their Official Mandate")
    roles = [
        ("1. Reception / Desk Operator", "Front Desk & Public Intake", "Enters citizen complaints, auto-fills details from verification reports, prints 80mm thermal receipts, and dispatches WhatsApp alerts.", C_CYAN),
        ("2. Circle Incharge (AD / DSP)", "Circle Command & Scrutiny", "Reviews incoming complaints, checks legal jurisdiction, assigns Verification Officers (VO), and approves enquiry recommendations.", C_GOLD),
        ("3. Verification Officer (VO)", "Field Verification Unit (ASI / SI)", "Conducts preliminary citizen inquiry, verifies phone/bank transactions, and submits verification findings with 1 of 4 clear outcomes.", C_GREEN),
        ("4. Enquiry Officer (IO)", "Core Investigation (Inspector / SI)", "Manages formal enquiries, issues Section 160 CrPC summons, records daily case diaries (Zimni), and prepares final reports (CFR).", C_CYAN),
        ("5. Digital Forensic Specialist", "Forensic Lab Division", "Takes custody of seized digital devices, conducts forensic data extraction, and attaches forensic certificates to the case dossier.", C_GOLD),
        ("6. Headquarters / DG Admin", "National Command & Monitoring", "Monitors nationwide Circle performance, pendency metrics, transfer cases across circles, and reviews institutional audit logs.", C_GREEN)
    ]
    for idx, (r_name, r_sub, r_desc, col) in enumerate(roles):
        row, col_idx = idx // 3, idx % 3
        c_left, c_top = 0.8 + col_idx * 3.98, 1.65 + row * 2.7
        add_card(s3, c_left, c_top, 3.75, 2.45, r_name, col, r_sub)
        sb = s3.shapes.add_textbox(Inches(c_left + 0.2), Inches(c_top + 0.85), Inches(3.35), Inches(1.45))
        stf = sb.text_frame
        stf.word_wrap = True
        p = stf.paragraphs[0]
        p.text = r_desc
        p.font.name = "Segoe UI"
        p.font.size = Pt(11)
        p.font.color.rgb = C_WHITE

    # SLIDE 4: Front Desk Workflow
    s4 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s4)
    add_header(s4, "Front Desk: Fast Registration & Receipt Issuance", "Step-by-Step Guide for Desk Operators and Reception Moharrars")
    add_card(s4, 0.8, 1.65, 5.75, 5.35, "📋 How Operator Registers a Complaint", C_CYAN)
    d_box = s4.shapes.add_textbox(Inches(1.05), Inches(2.3), Inches(5.25), Inches(4.5))
    dtf = d_box.text_frame
    dtf.word_wrap = True
    d_steps = [
        ("Step 1: Open Complaint Form", "Click '+ New Complaint' on the top right of the dashboard."),
        ("Step 2: Auto-Fill from Report", "If citizen brings a written verification report, drop the PDF into 'Import PDF' to auto-fill Name, CNIC, Phone, Address, and Accused Bank Details in 1 second."),
        ("Step 3: Verify Citizen Details", "Check Complainant CNIC (34100-XXXXXXX-X), Contact No, and Crime Type (Financial Fraud, Harassment, PPC 414, Identity Theft)."),
        ("Step 4: Save & Generate Tracking No", "System assigns unique Tracking Number (e.g. #0004/26).")
    ]
    for idx, (title, desc) in enumerate(d_steps):
        p = dtf.paragraphs[0] if idx == 0 else dtf.add_paragraph()
        p.text = f"{title}\n  ➔ {desc}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(8)

    add_card(s4, 6.78, 1.65, 5.75, 5.35, "🧾 80mm Thermal Receipt & WhatsApp Alert", C_GOLD)
    r_box = s4.shapes.add_textbox(Inches(7.03), Inches(2.3), Inches(5.25), Inches(4.5))
    rtf = r_box.text_frame
    rtf.word_wrap = True
    r_steps = [
        ("Official 80mm Thermal Slip", "Click 'Slip' button. Thermal printer instantly generates receipt featuring official NCCIA emblem, Complaint No, Date, and Verification Officer name."),
        ("Citizen WhatsApp Dispatch", "Click green 'WhatsApp' button. A pre-formatted official acknowledgment message with live tracking link is sent directly to the citizen's mobile."),
        ("QR Code Verification", "Citizen can scan the QR code anytime on their phone to see live progress without calling the station."),
        ("Zero Physical Register Entry", "Automates Moharrar manual diary register, eliminating handwritten errors.")
    ]
    for idx, (title, desc) in enumerate(r_steps):
        p = rtf.paragraphs[0] if idx == 0 else rtf.add_paragraph()
        p.text = f"{title}\n  ➔ {desc}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(8)

    # SLIDE 5: Incharge Scrutiny
    s5 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s5)
    add_header(s5, "Circle Incharge: Scrutiny & Officer Assignment", "How Assistant Directors / DSPs Supervise and Delegate Incoming Cases")
    add_card(s5, 0.8, 1.65, 5.75, 5.35, "⚖️ Legal Scrutiny of Complaints", C_GOLD)
    sc_box = s5.shapes.add_textbox(Inches(1.05), Inches(2.3), Inches(5.25), Inches(4.5))
    sctf = sc_box.text_frame
    sctf.word_wrap = True
    sc_items = [
        ("Review New Submissions", "Incharge dashboard shows all complaints registered in the Circle with real-time pending counters."),
        ("Check Legal Cognizance", "Verify if the offense falls under NCCIA cyber jurisdiction (PECA Act 2016 / Cyber Stalking / Banking Fraud)."),
        ("Mark Scrutiny Result", "Mark complaint as 'Valid & Proceed' or 'Irrelevant / Civil Dispute' with official remarks."),
        ("Set Priority Level", "Tag high-profile, financial embezzlement, or vulnerable victim cases as 'Urgent Priority'.")
    ]
    for idx, (title, desc) in enumerate(sc_items):
        p = sctf.paragraphs[0] if idx == 0 else sctf.add_paragraph()
        p.text = f"✔ {title}: {desc}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(10)

    add_card(s5, 6.78, 1.65, 5.75, 5.35, "👤 Direct Assignment to Verification Officer (VO)", C_CYAN)
    as_box = s5.shapes.add_textbox(Inches(7.03), Inches(2.3), Inches(5.25), Inches(4.5))
    astf = as_box.text_frame
    astf.word_wrap = True
    as_items = [
        ("One-Click Officer Marking", "Click 'Assign' next to the complaint. Select from active Circle Verification Officers (ASI / Sub-Inspector)."),
        ("Instant VO Notification", "Assigned officer immediately sees the new complaint in their pending investigation queue."),
        ("Automatic Assignment Slip Update", "The 80mm complaint slip and A4 report automatically reflect the assigned officer's name for public transparency."),
        ("Timeline Tracking", "Incharge monitors how many days a VO takes to submit the verification inquiry.")
    ]
    for idx, (title, desc) in enumerate(as_items):
        p = astf.paragraphs[0] if idx == 0 else astf.add_paragraph()
        p.text = f"✔ {title}: {desc}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(10)

    # SLIDE 6: VO Field Inquiry
    s6 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s6)
    add_header(s6, "Verification Officer: Field Inquiry & 4 Outcomes", "How VOs Conduct Preliminary Inquiries and Submit Official Findings")
    add_card(s6, 0.8, 1.65, 5.75, 5.35, "🔍 VO Investigation Procedure", C_GREEN)
    vo_box = s6.shapes.add_textbox(Inches(1.05), Inches(2.3), Inches(5.25), Inches(4.5))
    votf = vo_box.text_frame
    votf.word_wrap = True
    vo_steps = [
        ("1. Contact Complainant & Accused", "Call parties, record preliminary statements, and request supporting digital proof."),
        ("2. Verify Financial & CDR Trails", "Verify bank transaction IDs, mobile wallet transfers (Easypaisa/JazzCash), and IP logs."),
        ("3. Accused Profiling", "Add accused bank accounts, mobile numbers, and social media handles into the accused record table."),
        ("4. Prepare Verification Report", "Write summary of inquiry and select 1 of 4 official recommendations for Incharge approval.")
    ]
    for idx, (title, desc) in enumerate(vo_steps):
        p = votf.paragraphs[0] if idx == 0 else votf.add_paragraph()
        p.text = f"{title}\n  ➔ {desc}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(8)

    add_card(s6, 6.78, 1.65, 5.75, 5.35, "📊 4 Official Verification Recommendations", C_GOLD)
    oc_box = s6.shapes.add_textbox(Inches(7.03), Inches(2.3), Inches(5.25), Inches(4.5))
    octf = oc_box.text_frame
    octf.word_wrap = True
    outcomes = [
        ("1. Register Enquiry (Recommended)", "Digital evidence establishes prima facie offense. Case moves to Formal Enquiry Stage."),
        ("2. Case Closure", "Closed due to lack of evidence, non-pursuance by complainant, or mutual reconciliation."),
        ("3. Merge with Primary Complaint", "Accused is part of an ongoing syndicate already under active investigation in another complaint."),
        ("4. Transfer to Other Circle / Agency", "Offense occurred in another territorial Circle or belongs to Provincial Police / FIA / NAB.")
    ]
    for idx, (title, desc) in enumerate(outcomes):
        p = octf.paragraphs[0] if idx == 0 else octf.add_paragraph()
        p.text = f"🔸 {title}\n    {desc}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(8)

    # SLIDE 7: IO Enquiry & Case Diary
    s7 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s7)
    add_header(s7, "Investigation Officer (IO): Formal Enquiry & Case Diaries", "Managing Registered Enquiries, Section 160 CrPC Summons & Case Diaries (Zimni)")
    add_card(s7, 0.8, 1.65, 5.75, 5.35, "📂 Formal Enquiry Registration", C_CYAN)
    enq_box = s7.shapes.add_textbox(Inches(1.05), Inches(2.3), Inches(5.25), Inches(4.5))
    enqtf = enq_box.text_frame
    enqtf.word_wrap = True
    enq_items = [
        ("Enquiry Numbering", "System assigns official Enquiry ID (e.g. #LHR-E-80519/26)."),
        ("Assigned to Inspector / IO", "IO takes complete charge of the statutory investigation dossier."),
        ("Bank & Telecom Directives", "IO records formal requests sent to State Bank, Commercial Banks, and Telcos for account freezing and CDR."),
        ("Evidence Vault", "Uploads seized bank receipts, chat screenshots, and victim affidavits directly to the enquiry file.")
    ]
    for idx, (title, desc) in enumerate(enq_items):
        p = enqtf.paragraphs[0] if idx == 0 else enqtf.add_paragraph()
        p.text = f"✔ {title}: {desc}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(10)

    add_card(s7, 6.78, 1.65, 5.75, 5.35, "📖 Digital Case Diary (Zimni)", C_GOLD)
    zim_box = s7.shapes.add_textbox(Inches(7.03), Inches(2.3), Inches(5.25), Inches(4.5))
    zimtf = zim_box.text_frame
    zimtf.word_wrap = True
    zim_items = [
        ("Daily Progress Entries", "IO enters chronological Zimni entries with date, diary number, and detailed investigation notes."),
        ("Tamper-Proof Timestamps", "System locks diary timestamps, preventing retroactive tampering and ensuring court admissibility."),
        ("A4 Printable Zimni Dossier", "1-click generation of complete case diary history for submission to Circle Incharge or High Court."),
        ("Final Case Report (CFR)", "IO compiles Case Final Report with findings for FIR registration.")
    ]
    for idx, (title, desc) in enumerate(zim_items):
        p = zimtf.paragraphs[0] if idx == 0 else zimtf.add_paragraph()
        p.text = f"✔ {title}: {desc}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(10)

    # SLIDE 8: Notices Generator
    s8 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s8)
    add_header(s8, "Legal Notices Generator: Section 160 CrPC Summons", "Instant Generation of Statutory Appearance Notices in Thermal & A4 Formats")
    add_card(s8, 0.8, 1.65, 5.75, 5.35, "📜 Statutory Notice Features", C_GOLD)
    not_box = s8.shapes.add_textbox(Inches(1.05), Inches(2.3), Inches(5.25), Inches(4.5))
    nottf = not_box.text_frame
    nottf.word_wrap = True
    not_items = [
        ("Section 160 CrPC Compliance", "Generates legally binding call-up notices for accused and material witnesses."),
        ("Automated Officer & Station Details", "Auto-populates Circle Name, IO Name, Station Address, Date & Time of required appearance."),
        ("Anti-Fraud QR Verification", "Includes secure QR verification URL so recipient can verify the notice is genuine and issued by NCCIA."),
        ("Service Tracking", "Records delivery method (By Hand / Registered Post / WhatsApp / Courier) and recipient acknowledgment.")
    ]
    for idx, (title, desc) in enumerate(not_items):
        p = nottf.paragraphs[0] if idx == 0 else nottf.add_paragraph()
        p.text = f"🔹 {title}: {desc}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(10)

    add_card(s8, 6.78, 1.65, 5.75, 5.35, "🖨️ Dual Print Modes (80mm & A4)", C_CYAN)
    pm_box = s8.shapes.add_textbox(Inches(7.03), Inches(2.3), Inches(5.25), Inches(4.5))
    pmtf = pm_box.text_frame
    pmtf.word_wrap = True
    pm_items = [
        ("80mm Thermal Field Notice", "Compact, high-speed thermal notice printed instantly at the desk for hand-to-hand delivery by field officers."),
        ("A4 Formal Legal Notice", "Full official letterhead notice complete with NCCIA seal, legal warning, signature boxes, and duplicate copy for case file."),
        ("WhatsApp Electronic Notice", "Direct transmission of official summons to recipient's mobile with automatic timestamp logging."),
        ("Eliminates Notice Drafting Time", "Reduces notice creation from 15 minutes of typing to 1 single click.")
    ]
    for idx, (title, desc) in enumerate(pm_items):
        p = pmtf.paragraphs[0] if idx == 0 else pmtf.add_paragraph()
        p.text = f"🔹 {title}: {desc}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(10)

    # SLIDE 9: Digital Forensic Unit
    s9 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s9)
    add_header(s9, "Digital Forensic Unit: Evidence Chain of Custody", "Segregated Lab Portal for Mobile, Computer & Cloud Evidence Examination")
    add_card(s9, 0.8, 1.65, 5.75, 5.35, "🔬 Evidence Intake & Chain of Custody", C_CYAN)
    lab_box = s9.shapes.add_textbox(Inches(1.05), Inches(2.3), Inches(5.25), Inches(4.5))
    labtf = lab_box.text_frame
    labtf.word_wrap = True
    lab_items = [
        ("Seized Device Intake", "IO submits seized phones, laptops, USB drives, or DVRs with parcel and seal details."),
        ("Cryptographic Hash Calculation", "Records MD5 and SHA256 hashes upon intake to ensure zero evidence alteration."),
        ("Segregated Forensic Portal", "Forensic examiners access a dedicated /forensic interface without interfering with IO case files."),
        ("Hardware & OS Metadata", "Records IMEI, Serial Numbers, OS Version, and device condition.")
    ]
    for idx, (title, desc) in enumerate(lab_items):
        p = labtf.paragraphs[0] if idx == 0 else labtf.add_paragraph()
        p.text = f"✔ {title}: {desc}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(10)

    add_card(s9, 6.78, 1.65, 5.75, 5.35, "📑 Forensic Examination Report", C_GOLD)
    rep_box = s9.shapes.add_textbox(Inches(7.03), Inches(2.3), Inches(5.25), Inches(4.5))
    reptf = rep_box.text_frame
    reptf.word_wrap = True
    rep_items = [
        ("Extraction Findings", "Forensic specialist documents recovered WhatsApp chats, deleted call logs, crypto transactions, and financial artifacts."),
        ("Digital Certificate Attachment", "System generates official Forensic Report signed with examiner's digital signature."),
        ("Direct Attachment to Case File", "Report automatically links to the primary Enquiry/FIR file for court presentation."),
        ("Strict Chain-of-Custody Log", "Records who handled the evidence, date, time, and room storage location.")
    ]
    for idx, (title, desc) in enumerate(rep_items):
        p = reptf.paragraphs[0] if idx == 0 else reptf.add_paragraph()
        p.text = f"✔ {title}: {desc}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(10)

    # SLIDE 10: Court Trial
    s10 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s10)
    add_header(s10, "Case File (FIR) & Court Prosecution Tracking", "From FIR Registration and Challan Submission to Final Court Verdict")
    add_card(s10, 0.8, 1.65, 5.75, 5.35, "⚖️ FIR Registration & Challan", C_GOLD)
    fir_box = s10.shapes.add_textbox(Inches(1.05), Inches(2.3), Inches(5.25), Inches(4.5))
    firtf = fir_box.text_frame
    firtf.word_wrap = True
    fir_items = [
        ("FIR Creation", "Upon approval of CFR by Incharge/DG, system registers formal Case File with FIR Number (e.g. FIR # 12/26)."),
        ("Challan Preparation (Section 173 CrPC)", "Prosecution branch tracks Challan drafting, scrutiny, and submission date to Special Cyber Court."),
        ("Accused In Custody / Bail", "Tracks status of arrested suspects (Physical Remand, Judicial Lockup, Interim Bail, Proclaimed Offender)."),
        ("Seized Assets & Case Property", "Maintains register of frozen bank funds, recovered cash, and impounded vehicles.")
    ]
    for idx, (title, desc) in enumerate(fir_items):
        p = firtf.paragraphs[0] if idx == 0 else firtf.add_paragraph()
        p.text = f"🔸 {title}: {desc}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(10)

    add_card(s10, 6.78, 1.65, 5.75, 5.35, "🏛️ Court Hearings & Final Verdicts", C_GREEN)
    crt_box = s10.shapes.add_textbox(Inches(7.03), Inches(2.3), Inches(5.25), Inches(4.5))
    crttf = crt_box.text_frame
    crttf.word_wrap = True
    crt_items = [
        ("Hearing Calendar", "Automated alert for upcoming court appearance dates, IO witness cross-examinations, and prosecutor notes."),
        ("Bail Application Tracking", "Records Pre-Arrest and Post-Arrest bail decisions with judge orders."),
        ("Final Verdict Recording", "Records final court judgment: Conviction (Prison Term & Fines), Acquittal, or Compromise."),
        ("Recovery & Restitution", "Logs refunded amounts returned to scam victims through court directives.")
    ]
    for idx, (title, desc) in enumerate(crt_items):
        p = crttf.paragraphs[0] if idx == 0 else crttf.add_paragraph()
        p.text = f"🔸 {title}: {desc}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(10)

    # SLIDE 11: 10 Lakh Records Search
    s11 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s11)
    add_header(s11, "10 Lakh Historical Case Records Reference", "Instant Cross-Search Across Millions of Past FIRs, Enquiries & Scammer Accounts")
    add_card(s11, 0.8, 1.65, 5.75, 5.35, "🔍 Instant Multi-Parameter Search", C_CYAN)
    sr_box = s11.shapes.add_textbox(Inches(1.05), Inches(2.3), Inches(5.25), Inches(4.5))
    srtf = sr_box.text_frame
    srtf.word_wrap = True
    sr_items = [
        ("Search by CNIC", "Enter suspect or complainant CNIC to instantly pull their complete history across all Pakistan Circles."),
        ("Search by Bank Account / IBAN", "Detect if a scammer's bank account was involved in previous cyber fraud complaints in Karachi, Lahore, or Islamabad."),
        ("Search by Phone / WhatsApp", "Check call history, reported WhatsApp numbers, and suspect aliases."),
        ("Search by Historical Tracking ID", "Retrieve old legacy scanned files from the 10 Lakh historical database in milliseconds.")
    ]
    for idx, (title, desc) in enumerate(sr_items):
        p = srtf.paragraphs[0] if idx == 0 else srtf.add_paragraph()
        p.text = f"✔ {title}: {desc}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(10)

    add_card(s11, 6.78, 1.65, 5.75, 5.35, "📁 Fast Historical Ingestion Pipeline", C_GOLD)
    hi_box = s11.shapes.add_textbox(Inches(7.03), Inches(2.3), Inches(5.25), Inches(4.5))
    hitf = hi_box.text_frame
    hitf.word_wrap = True
    hi_items = [
        ("Hierarchical Folder Structure", "Historical files structured by Circle ➔ IO ➔ Moharar are automatically cataloged."),
        ("Zero Storage Duplication", "In-place referencing preserves server disk space while making 10 Lakh files searchable."),
        ("Historical Verification Reports", "Officers can open and read original legacy PDF reports with 1 click from their desk."),
        ("Syndicate Detection", "Cross-references past and present accused profiles to expose repeat cyber crime networks.")
    ]
    for idx, (title, desc) in enumerate(hi_items):
        p = hitf.paragraphs[0] if idx == 0 else hitf.add_paragraph()
        p.text = f"✔ {title}: {desc}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(10)

    # SLIDE 12: Circle Monitoring & Accountability
    s12 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s12)
    add_header(s12, "Circle Performance Monitoring & Accountability", "Real-Time Executive Dashboards for Incharges, Regional Directors & Headquarters")
    add_card(s12, 0.8, 1.65, 5.75, 5.35, "📊 Real-Time Circle KPIs", C_GOLD)
    kp_box = s12.shapes.add_textbox(Inches(1.05), Inches(2.3), Inches(5.25), Inches(4.5))
    kptf = kp_box.text_frame
    kptf.word_wrap = True
    kp_items = [
        ("Total Intake vs Disposal Rate", "Visual representation of complaints registered vs inquiries finalized per month."),
        ("Average Resolution Days", "Tracks average days taken by each Circle and Officer to complete investigations."),
        ("Crime Category Distribution", "Live breakdown of Financial Frauds, Harassment, Identity Theft, and PPC 414 cases."),
        ("Recovery Metrics", "Total defrauded money recovered and returned to citizens.")
    ]
    for idx, (title, desc) in enumerate(kp_items):
        p = kptf.paragraphs[0] if idx == 0 else kptf.add_paragraph()
        p.text = f"📈 {title}: {desc}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(10)

    add_card(s12, 6.78, 1.65, 5.75, 5.35, "🛡️ Officer Login & Security Audit Trail", C_CYAN)
    au_box = s12.shapes.add_textbox(Inches(7.03), Inches(2.3), Inches(5.25), Inches(4.5))
    autf = au_box.text_frame
    autf.word_wrap = True
    au_items = [
        ("Real-Time IP Logging", "Tracks connecting IP, Real Origin IP, and login timestamps for every officer session."),
        ("Spoof & Unauthorized Access Alerts", "Flags any suspicious proxy or VPN usage instantly on the security dashboard."),
        ("Action History Audit", "Every status update, case assignment, report submission, and notice print is recorded with user stamp."),
        ("Zero Unauthorized Alteration", "Ensures high institutional discipline and security compliance.")
    ]
    for idx, (title, desc) in enumerate(au_items):
        p = autf.paragraphs[0] if idx == 0 else autf.add_paragraph()
        p.text = f"🛡️ {title}: {desc}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(10)

    # SLIDE 13: Live Demo Walkthrough
    s13 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s13)
    add_header(s13, "Live Practical Demonstration Guide", "Step-by-Step Practical Demonstration for Officers & Leadership")
    live_steps = [
        ("1. Desk Operator Demo", "Open /complaints/create ➔ Auto-fill from sample report ➔ Register complaint #0004/26 ➔ Print 80mm thermal slip with VO name & QR ➔ Send WhatsApp alert."),
        ("2. Circle Incharge Demo", "Open Dashboard ➔ Review pending complaints ➔ Assign Verification Officer (VO) with High Priority tag."),
        ("3. VO Inquiry Demo", "VO reviews case evidence ➔ Records preliminary findings ➔ Selects Register Enquiry recommendation ➔ Submits report."),
        ("4. IO Enquiry & Notice Demo", "IO opens registered enquiry #LHR-E-80519/26 ➔ Generates Section 160 CrPC Legal Notice (80mm & A4) ➔ Enters daily Case Diary (Zimni)."),
        ("5. Forensic Lab Demo", "Open /forensic ➔ Review evidence item ➔ Verify hash ➔ Attach Forensic Examination Certificate."),
        ("6. Security & Audit Demo", "Open /login-history ➔ Search IP address ➔ Demonstrate live audit trail and real IP verification.")
    ]
    for idx, (title, desc) in enumerate(live_steps):
        row, col_idx = idx // 2, idx % 2
        c_left, c_top = 0.8 + col_idx * 5.98, 1.65 + row * 1.8
        add_card(s13, c_left, c_top, 5.75, 1.62, title, C_GOLD if idx == 0 else C_CYAN)
        sb = s13.shapes.add_textbox(Inches(c_left + 0.2), Inches(c_top + 0.5), Inches(5.35), Inches(1.05))
        stf = sb.text_frame
        stf.word_wrap = True
        p = stf.paragraphs[0]
        p.text = desc
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_WHITE

    # SLIDE 14: SOPs & Conclusion
    s14 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s14)
    add_header(s14, "Standard Operating Procedures (SOPs) & Next Steps", "Guidelines for Full Agency Rollout Across All Circles Nationwide")
    add_card(s14, 0.8, 1.65, 11.733, 5.35, "🎖️ Standards of Official Operational Excellence", C_GOLD)
    fn_box = s14.shapes.add_textbox(Inches(1.1), Inches(2.35), Inches(11.133), Inches(4.4))
    fntf = fn_box.text_frame
    fntf.word_wrap = True
    sop_points = [
        "Mandatory 80mm Receipt Issuance: Every complainant visiting an NCCIA station must receive a printed 80mm thermal receipt with tracking QR code.",
        "24-Hour Scrutiny Turnaround: Circle Incharges must review and assign complaints within 24 hours of front desk registration.",
        "Daily Case Diary (Zimni) Compliance: Investigation Officers must enter daily progress in the digital diary to ensure strict court admissibility.",
        "Zero Paper Dossier Mandate: All bank correspondence, evidence items, and notices must be uploaded to the central digital dossier.",
        "Accountability & Citizen Service: Direct WhatsApp tracking and transparency empower citizens and strengthen public confidence in NCCIA."
    ]
    for idx, pt in enumerate(sop_points):
        p = fntf.paragraphs[0] if idx == 0 else fntf.add_paragraph()
        p.text = f"⭐ {pt}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(12)

    output_path = os.path.abspath("f:/NCCIA/NCCIA_Officer_Operational_CMS_Presentation.pptx")
    prs.save(output_path)
    print(f"Successfully generated: {output_path}")

if __name__ == "__main__":
    create_officer_presentation()
