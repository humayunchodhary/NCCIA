import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    C_NAVY_BG      = RGBColor(11, 22, 44)
    C_CARD_BG      = RGBColor(19, 36, 68)
    C_CARD_BORDER  = RGBColor(41, 74, 110)
    C_GOLD         = RGBColor(229, 185, 83)
    C_GREEN        = RGBColor(16, 149, 93)
    C_CYAN         = RGBColor(56, 189, 248)
    C_WHITE        = RGBColor(255, 255, 255)
    C_SLATE        = RGBColor(160, 174, 192)

    logo_path = os.path.abspath('f:/NCCIA/public/images/images.jpg')

    def add_slide_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_NAVY_BG
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, subtitle_text=''):
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(10.5), Inches(1.15))
        tf = header_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Georgia'
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = C_GOLD
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = 'Segoe UI'
            p2.font.size = Pt(12.5)
            p2.font.bold = True
            p2.font.color.rgb = C_CYAN
            p2.space_before = Pt(4)
            
        if os.path.isfile(logo_path):
            slide.shapes.add_picture(logo_path, Inches(12.0), Inches(0.3), height=Inches(1.05))

    def add_card(slide, left, top, width, height, title='', title_color=C_GOLD, badge=''):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD_BG
        card.line.color.rgb = C_CARD_BORDER
        card.line.width = Pt(1.5)
        
        if title:
            tb = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.18), Inches(width - 0.5), Inches(0.45))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = 'Segoe UI'
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = title_color
            
            if badge:
                p2 = tf.add_paragraph()
                p2.text = badge
                p2.font.name = 'Segoe UI'
                p2.font.size = Pt(10)
                p2.font.bold = True
                p2.font.color.rgb = C_GREEN
        return card

    # SLIDE 1: Title
    s1 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s1)
    logo_size = 2.4
    logo_left = (13.333 - logo_size) / 2
    logo_top = 0.8
    ring = s1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(logo_left - 0.08), Inches(logo_top - 0.08), Inches(logo_size + 0.16), Inches(logo_size + 0.16))
    ring.fill.solid()
    ring.fill.fore_color.rgb = C_GOLD
    ring.line.fill.background()

    if os.path.isfile(logo_path):
        s1.shapes.add_picture(logo_path, Inches(logo_left), Inches(logo_top), width=Inches(logo_size), height=Inches(logo_size))

    tbox = s1.shapes.add_textbox(Inches(1.0), Inches(3.4), Inches(11.333), Inches(3.7))
    tf1 = tbox.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    p1.text = 'NATIONAL CYBER CRIME INVESTIGATION AGENCY'
    p1.font.name = 'Georgia'
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = C_GOLD

    p2 = tf1.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = 'قومی سائبر کرائم تفتیش ایجنسی — وزارت داخلہ حکومت پاکستان'
    p2.font.name = 'Segoe UI'
    p2.font.size = Pt(15)
    p2.font.bold = True
    p2.font.color.rgb = C_WHITE
    p2.space_before = Pt(4)

    p3 = tf1.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.text = 'Case Management System (CMS) — Officer Training & Operational SOPs'
    p3.font.name = 'Segoe UI'
    p3.font.size = Pt(20)
    p3.font.bold = True
    p3.font.color.rgb = C_CYAN
    p3.space_before = Pt(10)

    p4 = tf1.add_paragraph()
    p4.alignment = PP_ALIGN.CENTER
    p4.text = 'A Complete Practical Guide for Desk Operators, Circle Incharges, VOs, IOs & Lab Experts'
    p4.font.name = 'Segoe UI'
    p4.font.size = Pt(13)
    p4.font.color.rgb = C_SLATE
    p4.space_before = Pt(6)

    # SLIDE 2: Purpose
    s2 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s2)
    add_header(s2, 'Purpose & Vision of the CMS Portal', 'Transforming Traditional Paper Policing into Transparent Digital Workflows')
    p_cards = [
        ('📁 100% Digital Case Record', 'Eliminates physical file loss or delay. Every complaint, evidence, notice, and diary entry is permanently secured in the central digital database.', C_CYAN),
        ('⏱️ Fast Citizen Registration', 'Front desk registration takes less than 60 seconds with instant 80mm thermal slip generation and automated WhatsApp alert dispatch.', C_GOLD),
        ('⚖️ Officer Accountability', 'Complete timeline tracking from initial filing to scrutiny, verification, enquiry, and court trial. Zero hidden pendency.', C_GREEN),
        ('🔍 Instant Case Search', 'Find any past or present case instantly by Complaint Number, Complainant CNIC, Phone Number, Accused Bank Account, or Offence Category.', C_CYAN)
    ]
    for idx, (title, desc, col) in enumerate(p_cards):
        row, col_idx = idx // 2, idx % 2
        c_left, c_top = 0.8 + col_idx * 5.98, 1.65 + row * 2.7
        add_card(s2, c_left, c_top, 5.75, 2.45, title, col)
        sb = s2.shapes.add_textbox(Inches(c_left + 0.25), Inches(c_top + 0.65), Inches(5.25), Inches(1.65))
        stf = sb.text_frame
        stf.word_wrap = True
        p = stf.paragraphs[0]
        p.text = desc
        p.font.name = 'Segoe UI'
        p.font.size = Pt(12)
        p.font.color.rgb = C_WHITE
