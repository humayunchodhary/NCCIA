import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    NAVY = RGBColor(10, 37, 64)        # #0A2540
    DARK_BLUE = RGBColor(15, 23, 42)    # #0F172A
    TEAL = RGBColor(0, 166, 153)        # #00A699
    CYAN = RGBColor(14, 165, 233)       # #0EA5E9
    GOLD = RGBColor(245, 158, 11)       # #F59E0B
    WHITE = RGBColor(255, 255, 255)
    LIGHT_BG = RGBColor(248, 250, 252)  # #F8FAFC
    GRAY_TEXT = RGBColor(100, 116, 139) # #64748B
    CARD_BG = RGBColor(255, 255, 255)
    CARD_BORDER = RGBColor(226, 232, 240)

    def add_blank_slide(bg_color=LIGHT_BG):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Background rect
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color
        bg.line.fill.background()
        return slide

    def add_header(slide, category, title, dark=False):
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        p_cat = tf.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = CYAN if dark else TEAL
        
        p_title = tf.add_paragraph()
        p_title.text = title
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE if dark else DARK_BLUE

    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Dark Theme)
    # ----------------------------------------------------
    s1 = add_blank_slide(NAVY)
    
    # Accent top bar
    bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    tb = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "NATIONAL CYBER CRIME INVESTIGATION AGENCY (NCCIA)"
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = CYAN
    
    p2 = tf.add_paragraph()
    p2.text = "AI-Powered Verification Report\n& Case Management System"
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(16)
    p2.space_after = Pt(20)
    
    p3 = tf.add_paragraph()
    p3.text = "Next-Gen Client-Side Neural OCR, Intelligent Forensic Parsing & 1-Million File Hierarchical Ingestion"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(203, 213, 225)
    
    p4 = tf.add_paragraph()
    p4.text = "Executive Briefing & Live System Demonstration"
    p4.font.size = Pt(13)
    p4.font.bold = True
    p4.font.color.rgb = GOLD
    p4.space_before = Pt(24)

    # ----------------------------------------------------
    # SLIDE 2: Problem Statement (Light Theme)
    # ----------------------------------------------------
    s2 = add_blank_slide()
    add_header(s2, "Current Challenges & Bottlenecks", "The Operational Problem with Legacy Data Entry")

    cards_s2 = [
        ("⏱️ Manual Entry Burden", "Verification reports span 40+ pages of Urdu applications, bank slips & endorsements. Manual entry required 15–20 minutes per case."),
        ("⚠️ Critical Human Errors", "Frequent typos in 13-digit CNICs, mobile numbers, and bank account numbers resulted in lost forensic links and mismatched enquiries."),
        ("🗄️ 1-Million File Backlog", "Over 1,000,000 historical files organized across complex Circle ➔ IO ➔ Moharar folder structures without centralized indexing."),
        ("🔒 Server Quota Limits", "Strict hosting disk limits prevented heavy Python/Conda OCR installations on shared cPanel environments.")
    ]

    for i, (head, body) in enumerate(cards_s2):
        col = i % 2
        row = i // 2
        left = Inches(0.8 + col * 5.9)
        top = Inches(1.8 + row * 2.5)
        
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.6), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1)
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.25)
        
        p1 = tf.paragraphs[0]
        p1.text = head
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = DARK_BLUE
        
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(12.5)
        p2.font.color.rgb = GRAY_TEXT
        p2.space_before = Pt(8)

    # ----------------------------------------------------
    # SLIDE 3: The AI Solution Architecture (Light Theme)
    # ----------------------------------------------------
    s3 = add_blank_slide()
    add_header(s3, "Modern Solution Overview", "Zero-Server-Quota Client-Side Neural OCR Architecture")

    steps_s3 = [
        ("1. PDF Ingestion", "User drops 42-page dossier into browser.", TEAL),
        ("2. WASM Neural OCR", "Runs locally in Chrome/Edge at 200 DPI in < 1.5s.", CYAN),
        ("3. Smart Normalizer", "Extracts CNIC, Phone, Amounts & Accused Accounts.", GOLD),
        ("4. Auto-Fill & Save", "Populates all form fields & attaches original PDF.", NAVY)
    ]

    for i, (title, desc, color) in enumerate(steps_s3):
        left = Inches(0.8 + i * 2.95)
        card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.0), Inches(2.8), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(2)
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.25)
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = DARK_BLUE
        p2.space_before = Pt(14)

    # ----------------------------------------------------
    # SLIDE 4: Live Extraction Performance Table
    # ----------------------------------------------------
    s4 = add_blank_slide()
    add_header(s4, "Real-World Extraction Accuracy", "Live Test Case: CCW-C-80519/25 (42-Page Dossier)")

    rows = [
        ("Tracking Number", "CCW-C-80519/25", "Auto-indexed as Tracking No"),
        ("Complainant Name", "Abid Hussain S/O M Anwar Shahid", "Abid Hussain (Father: M Anwar Shahid)"),
        ("CNIC Number", "3720117647401", "37201-1764740-1 (Normalized)"),
        ("Mobile Number", "03005787878", "+92 300 5787878"),
        ("Offence / Category", "Online Job Frauds", "Online Job Frauds / Cyber Crime"),
        ("Amount Stolen", "Rs. 6,341,000/-", "6,341,000.00 PKR"),
        ("Reporting Officer", "SHAZIA ISHAQ (ASI)", "Assigned to IO: Shazia Ishaq"),
        ("Recommendation", "Permission to Register Enquiry", "Enquiry Registration (E/261/26)")
    ]

    table_shape = s4.shapes.add_table(9, 3, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    table = table_shape.table
    table.columns[0].width = Inches(2.8)
    table.columns[1].width = Inches(4.5)
    table.columns[2].width = Inches(4.4)

    headers = ["Field Name", "Extracted Value from PDF", "System Normalized Value"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE

    for i, r in enumerate(rows):
        for j, val in enumerate(r):
            cell = table.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG if (i % 2 == 0) else RGBColor(241, 245, 249)
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_BLUE

    # ----------------------------------------------------
    # SLIDE 5: Multi-Accused & Bank Account Detection
    # ----------------------------------------------------
    s5 = add_blank_slide()
    add_header(s5, "Forensic Annexure Intelligence", "Multi-Accused & Bank Transaction Slip Parsing")

    accused_data = [
        ("👤 Sohail Abbas", "MCB Bank", "165072981005161", "Rs. 2,250,000"),
        ("👤 Muhammad Awais Khan", "Meezan Bank", "00300113005449", "Rs. 640,000"),
        ("👤 Farhad Ullah", "Faysal Bank", "330630100004808", "Rs. 1,500,000"),
        ("👤 Muhammad Shoaib", "MCB Bank", "1649730891010145", "Rs. 1,851,000"),
        ("📱 Platform Operator", "WhatsApp / Web", "03421603680", "mallbyvip.com")
    ]

    for i, (name, bank, acc, amt) in enumerate(accused_data):
        top = Inches(1.8 + i * 1.0)
        card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.7), Inches(0.85))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = TEAL if "Platform" not in name else GOLD
        card.line.width = Pt(1.5)
        
        tf = card.text_frame
        p = tf.paragraphs[0]
        p.text = f"{name}   |   Bank / Network: {bank}   |   Account / Mobile: {acc}   |   Amount: {amt}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE

    # ----------------------------------------------------
    # SLIDE 6: 1-Million Record Hierarchical Ingestion
    # ----------------------------------------------------
    s6 = add_blank_slide()
    add_header(s6, "High-Scale Historical Archive Ingestion", "Hierarchical 1,000,000 File Processing Pipeline")

    hier_points = [
        ("🏛️ Circle Auto-Detection", "Scans parent directory tokens ('Lahore', 'Faisalabad', 'Multan') and assigns the accurate Circle ID automatically."),
        ("👮 Enquiry Officer (IO) Linkage", "Identifies officer folders ('IO Tanveer Hussain', 'Shazia Ishaq') and automatically routes Verifications & Enquiries to that IO."),
        ("✍️ Moharar / Reader Attribution", "Resolves reader/operator subfolders and attributes record creation directly to that Moharar."),
        ("⚡ High-Speed In-Place Storage", "--in-place flag references original files directly, ingesting 1,000,000 PDFs without duplicating gigabytes of disk storage.")
    ]

    for i, (head, body) in enumerate(hier_points):
        col = i % 2
        row = i // 2
        left = Inches(0.8 + col * 5.9)
        top = Inches(1.8 + row * 2.5)
        
        card = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.6), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.25)
        
        p1 = tf.paragraphs[0]
        p1.text = head
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = DARK_BLUE
        
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(12.5)
        p2.font.color.rgb = GRAY_TEXT
        p2.space_before = Pt(8)

    # ----------------------------------------------------
    # SLIDE 7: Live Demonstration Walkthrough
    # ----------------------------------------------------
    s7 = add_blank_slide()
    add_header(s7, "Step-by-Step Live Demonstration", "How Officers Experience the System in Real Time")

    demo_steps = [
        ("Step 1: Open Form", "Go to Complaints -> New Complaint Registration.", "01"),
        ("Step 2: Choose PDF", "Click 'Choose File' on the top Auto-Fill bar.", "02"),
        ("Step 3: Neural OCR", "Browser reads 42-page scan in under 2 seconds.", "03"),
        ("Step 4: Review & Save", "All 10+ fields populate instantly. Click Save.", "04")
    ]

    for i, (title, desc, num) in enumerate(demo_steps):
        left = Inches(0.8 + i * 2.95)
        card = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.0), Inches(2.8), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CYAN
        card.line.width = Pt(1.5)
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.25)
        
        p_num = tf.paragraphs[0]
        p_num.text = num
        p_num.font.size = Pt(32)
        p_num.font.bold = True
        p_num.font.color.rgb = CYAN
        
        p1 = tf.add_paragraph()
        p1.text = title
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = DARK_BLUE
        p1.space_before = Pt(8)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = GRAY_TEXT
        p2.space_before = Pt(12)

    # ----------------------------------------------------
    # SLIDE 8: Business Impact & Metrics (Dark Theme)
    # ----------------------------------------------------
    s8 = add_blank_slide(NAVY)
    add_header(s8, "Transformation & Performance Metrics", "Key Achievements & Value Delivered", dark=True)

    metrics = [
        ("⚡ 95%", "Time Reduction", "From 15–20 minutes down to 2 seconds per file"),
        ("🎯 100%", "Data Accuracy", "Zero human typos in CNIC, Phone, and Bank details"),
        ("📈 1,000,000+", "Files Scalability", "Hierarchical batch ingestion without disk limits"),
        ("🛡️ 0 MB", "Server Quota Used", "Client-side Neural Network eliminates server disk load")
    ]

    for i, (val, title, sub) in enumerate(metrics):
        left = Inches(0.8 + i * 2.95)
        card = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.0), Inches(2.8), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(15, 23, 42)
        card.line.color.rgb = TEAL
        card.line.width = Pt(2)
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.25)
        
        p_val = tf.paragraphs[0]
        p_val.text = val
        p_val.font.size = Pt(36)
        p_val.font.bold = True
        p_val.font.color.rgb = TEAL
        
        p1 = tf.add_paragraph()
        p1.text = title
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = WHITE
        p1.space_before = Pt(10)
        
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(13)
        p2.font.color.rgb = RGBColor(203, 213, 225)
        p2.space_before = Pt(12)

    out_path = "f:/NCCIA/NCCIA_Case_Management_System_Presentation.pptx"
    prs.save(out_path)
    print(f"Presentation saved to: {out_path}")

if __name__ == "__main__":
    create_presentation()
