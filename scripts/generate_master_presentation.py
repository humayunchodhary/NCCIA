import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_master_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette
    C_NAVY_BG      = RGBColor(10, 25, 47)      # #0A192F
    C_CARD_BG      = RGBColor(17, 34, 64)      # #112240
    C_CARD_BORDER  = RGBColor(35, 53, 84)      # #233554
    C_GOLD         = RGBColor(212, 175, 55)    # #D4AF37
    C_CYAN         = RGBColor(0, 180, 216)     # #00B4D8
    C_WHITE        = RGBColor(255, 255, 255)
    C_SLATE        = RGBColor(148, 163, 184)   # #94A3B8
    C_GREEN        = RGBColor(16, 185, 129)    # #10B981

    logo_path = os.path.abspath("f:/NCCIA/public/images/images.jpg")

    def add_slide_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_NAVY_BG
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, subtitle_text=""):
        # Header bar
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10.5), Inches(1.1))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Georgia"
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = C_GOLD
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = "Segoe UI"
            p2.font.size = Pt(13)
            p2.font.color.rgb = C_CYAN
            p2.space_before = Pt(4)
            
        # Top right small logo
        if os.path.isfile(logo_path):
            slide.shapes.add_picture(logo_path, Inches(12.0), Inches(0.35), height=Inches(1.0))

    def add_card(slide, left, top, width, height, title="", title_color=C_CYAN):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD_BG
        card.line.color.rgb = C_CARD_BORDER
        card.line.width = Pt(1.5)
        
        if title:
            tb = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.2), Inches(width - 0.5), Inches(0.45))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = "Segoe UI"
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = title_color
        return card

    # ==========================================
    # SLIDE 1: ULTRA-LUXURY TITLE SLIDE
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s1)

    # Decorative Gold Ring for Logo
    logo_size = 2.4
    logo_left = (13.333 - logo_size) / 2
    logo_top = 0.85
    
    ring = s1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(logo_left - 0.08), Inches(logo_top - 0.08), Inches(logo_size + 0.16), Inches(logo_size + 0.16))
    ring.fill.solid()
    ring.fill.fore_color.rgb = C_GOLD
    ring.line.fill.background()

    if os.path.isfile(logo_path):
        s1.shapes.add_picture(logo_path, Inches(logo_left), Inches(logo_top), width=Inches(logo_size), height=Inches(logo_size))

    # Title & Subtitle Box
    tbox = s1.shapes.add_textbox(Inches(1.0), Inches(3.45), Inches(11.333), Inches(3.6))
    tf1 = tbox.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    p1.text = "NATIONAL CYBER CRIME INVESTIGATION AGENCY"
    p1.font.name = "Georgia"
    p1.font.size = Pt(30)
    p1.font.bold = True
    p1.font.color.rgb = C_GOLD

    p2 = tf1.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = "قومی سائبر کرائم تفتیش ایجنسی — حکومت پاکستان"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = C_WHITE
    p2.space_before = Pt(6)

    p3 = tf1.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.text = "Enterprise Digital Case Management & Forensic Intelligence System"
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(20)
    p3.font.bold = True
    p3.font.color.rgb = C_CYAN
    p3.space_before = Pt(12)

    p4 = tf1.add_paragraph()
    p4.alignment = PP_ALIGN.CENTER
    p4.text = "High-Concurrency Architecture (10M+ Ready) | Client-Side AI WASM OCR | End-to-End Investigation Lifecycle"
    p4.font.name = "Segoe UI"
    p4.font.size = Pt(13)
    p4.font.color.rgb = C_SLATE
    p4.space_before = Pt(8)

    # ==========================================
    # SLIDE 2: EXECUTIVE SUMMARY & STATS
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s2)
    add_header(s2, "Executive Overview & Core Pillars", "Modernizing National Cyber Crime Investigation with Next-Gen Intelligence")

    # 4 Key Stat Callouts
    stats_data = [
        ("1-Second", "AI Neural OCR", "Zero server CPU load with browser-native WASM engine"),
        ("10 Lakh+", "Historical Ingestion", "Automated Circle ➔ IO ➔ Moharar folder ingestion"),
        ("10M+ Req", "Scale Optimized", "B-Tree composite indexed database & Redis caching"),
        ("100% Zero-Scroll", "Responsive Design", "Pixel-perfect experience across mobile, tablet, laptop")
    ]
    for idx, (val, label, sub) in enumerate(stats_data):
        c_left = 0.8 + idx * 2.98
        add_card(s2, c_left, 1.7, 2.8, 1.7)
        stb = s2.shapes.add_textbox(Inches(c_left + 0.15), Inches(1.85), Inches(2.5), Inches(1.4))
        stf = stb.text_frame
        stf.word_wrap = True
        
        sp1 = stf.paragraphs[0]
        sp1.text = val
        sp1.font.name = "Segoe UI"
        sp1.font.size = Pt(24)
        sp1.font.bold = True
        sp1.font.color.rgb = C_GOLD
        
        sp2 = stf.add_paragraph()
        sp2.text = label
        sp2.font.name = "Segoe UI"
        sp2.font.size = Pt(13)
        sp2.font.bold = True
        sp2.font.color.rgb = C_WHITE
        sp2.space_before = Pt(2)
        
        sp3 = stf.add_paragraph()
        sp3.text = sub
        sp3.font.name = "Segoe UI"
        sp3.font.size = Pt(10)
        sp3.font.color.rgb = C_SLATE
        sp3.space_before = Pt(4)

    # Bottom 2 Cards
    add_card(s2, 0.8, 3.65, 5.75, 3.3, "🛡️ National Security Mission", C_CYAN)
    m_box = s2.shapes.add_textbox(Inches(1.05), Inches(4.25), Inches(5.25), Inches(2.5))
    mtf = m_box.text_frame
    mtf.word_wrap = True
    mission_points = [
        "Digital Sovereign Platform: Complete digital transformation from physical paper records to high-security electronic custody.",
        "Rapid Public Redressal: Automated 80mm receipt generation & real-time WhatsApp case tracking for citizens.",
        "Forensic Integrity: Strict chain-of-custody tracking with cryptographic MD5/SHA256 evidence hashing."
    ]
    for idx, pt in enumerate(mission_points):
        p = mtf.paragraphs[0] if idx == 0 else mtf.add_paragraph()
        p.text = "• " + pt
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(8)

    add_card(s2, 6.78, 3.65, 5.75, 3.3, "⚡ Key Architectural Highlights", C_GOLD)
    a_box = s2.shapes.add_textbox(Inches(7.03), Inches(4.25), Inches(5.25), Inches(2.5))
    atf = a_box.text_frame
    atf.word_wrap = True
    arch_points = [
        "Zero-Lag Query Engine: Full-table scan elimination with multi-column composite indexing.",
        "Anti-Brute Force Security: 6-attempt throttle, real proxy IP extraction, and automated spoof detection.",
        "Unified Official Identity: Authentic circular NCCIA seal on all 80mm thermal slips, A4 legal reports, and case diaries."
    ]
    for idx, pt in enumerate(arch_points):
        p = atf.paragraphs[0] if idx == 0 else atf.add_paragraph()
        p.text = "• " + pt
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(8)

    # ==========================================
    # SLIDE 3: AI-POWERED ZERO-QUOTA OCR
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s3)
    add_header(s3, "AI-Powered PDF OCR Neural Engine", "Zero Server Storage Quota Overhead · 1-Second Auto-Fill")

    add_card(s3, 0.8, 1.7, 5.75, 5.25, "🧠 Browser-Native WebAssembly Engine", C_CYAN)
    w_box = s3.shapes.add_textbox(Inches(1.05), Inches(2.35), Inches(5.25), Inches(4.3))
    wtf = w_box.text_frame
    wtf.word_wrap = True
    wasm_points = [
        ("Client-Side Neural Execution", "Tesseract.js v7 running inside browser WebAssembly (WASM), rendering 200 DPI anti-aliased canvas frames on user GPU/CPU."),
        ("Zero Hosting Disk Quota", "Zero Python/Conda dependencies on hosting server — 100% immune to shared hosting user block limit errors."),
        ("Multi-Page Intelligent Parsing", "Scans 40+ page CCW verification PDFs, extracting Complainant Name, CNIC, Phone, Address, Offence, Amount, and Accused Bank Accounts in 1 second."),
        ("Direct JSON Hydration", "Instantly populates React state forms with clean formatting, auto-hyphenated CNIC, and structured accused bank tables.")
    ]
    for idx, (title, desc) in enumerate(wasm_points):
        p = wtf.paragraphs[0] if idx == 0 else wtf.add_paragraph()
        p.text = f"✔ {title}: {desc}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(10)

    add_card(s3, 6.78, 1.7, 5.75, 5.25, "📋 Auto-Extracted Intelligence Fields", C_GOLD)
    f_box = s3.shapes.add_textbox(Inches(7.03), Inches(2.35), Inches(5.25), Inches(4.3))
    ftf = f_box.text_frame
    ftf.word_wrap = True
    field_items = [
        ("Complainant Identity", "Name, Father Name, CNIC (34100-XXXXXXX-X), Mobile No, Address, District"),
        ("Case Classification", "Offence Category (Financial Fraud, Harassment, PPC 414), Diary No, Report Date"),
        ("Financial Trails", "Amount defrauded, Sender Bank & Account No, Receiver Bank & IBAN Account No"),
        ("Accused Profiles", "Accused Name, Father Name, Contact, Social Media URLs, and CNIC/Passport details"),
        ("Investigating Hierarchy", "Circle Name, Verification Officer (VO), and Enquiry Officer (IO) assignments")
    ]
    for idx, (title, desc) in enumerate(field_items):
        p = ftf.paragraphs[0] if idx == 0 else ftf.add_paragraph()
        p.text = f"🔹 {title}\n    {desc}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(10)

    # ==========================================
    # SLIDE 4: 10 LAKH HISTORICAL BULK INGESTION
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s4)
    add_header(s4, "10 Lakh Historical File Bulk Ingestion", "High-Throughput In-Place Archival · Circle ➔ IO ➔ Moharar Hierarchy")

    add_card(s4, 0.8, 1.7, 3.7, 5.25, "📁 Hierarchy Auto-Discovery", C_CYAN)
    h_box = s4.shapes.add_textbox(Inches(1.0), Inches(2.35), Inches(3.3), Inches(4.3))
    htf = h_box.text_frame
    htf.word_wrap = True
    h_items = [
        "Recursive Path Scanner: Automatically parses directory structures like `Lahore/AD_Tanveer/Moharar_Cases/*.pdf`.",
        "Entity Resolution: Maps folder names to Circles, Zones, and registered Police Officers.",
        "Smart Deduplication: Prevents double-processing via tracking number and MD5 checksum validation."
    ]
    for idx, item in enumerate(h_items):
        p = htf.paragraphs[0] if idx == 0 else htf.add_paragraph()
        p.text = "• " + item
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(12)

    add_card(s4, 4.8, 1.7, 3.7, 5.25, "⚡ In-Place Zero-Copy Engine", C_GOLD)
    i_box = s4.shapes.add_textbox(Inches(5.0), Inches(2.35), Inches(3.3), Inches(4.3))
    itf = i_box.text_frame
    itf.word_wrap = True
    i_items = [
        "`--in-place` Mode: References historical PDFs directly at their original network storage location without file duplication.",
        "Zero Storage Bloat: 10 Lakh files ingested without consuming extra disk gigabytes.",
        "Batch Transaction Pipeline: Processes 5,000 files/minute with safe rollbacks and recovery logs."
    ]
    for idx, item in enumerate(i_items):
        p = itf.paragraphs[0] if idx == 0 else itf.add_paragraph()
        p.text = "• " + item
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(12)

    add_card(s4, 8.8, 1.7, 3.7, 5.25, "📊 Progress & Error Telemetry", C_GREEN)
    p_box = s4.shapes.add_textbox(Inches(9.0), Inches(2.35), Inches(3.3), Inches(4.3))
    ptf = p_box.text_frame
    ptf.word_wrap = True
    p_items = [
        "Live CLI & Web Stats: Real-time progress bars, success counters, and failure diagnostic reports.",
        "Automatic OCR Fallback: High-speed text layer extraction with fallback to OCR for scanned physical papers.",
        "Audit Log Trail: Full database history recording ingestion timestamp, file size, and assigned officer."
    ]
    for idx, item in enumerate(p_items):
        p = ptf.paragraphs[0] if idx == 0 else ptf.add_paragraph()
        p.text = "• " + item
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(12)

    # ==========================================
    # SLIDE 5: VISUAL 6-STAGE INVESTIGATION WORKFLOW
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s5)
    add_header(s5, "End-to-End Investigation Lifecycle", "6 Seamless Stages from Complaint Filing to Court Verdict")

    stages = [
        ("Stage 1", "Registration", "Desk Operator / PDF Auto-Fill\n80mm Slip & WhatsApp Alert", C_CYAN),
        ("Stage 2", "Scrutiny", "Circle Incharge Legal Scrutiny\nAssign Verification Officer", C_GOLD),
        ("Stage 3", "Verification", "VO Field Inquiry & CDR Analysis\n4 Recommendation Matrix", C_GREEN),
        ("Stage 4", "Formal Enquiry", "IO Formal Enquiry Registration\nLegal Notices (80mm/A4 QR)", C_CYAN),
        ("Stage 5", "Forensic Lab", "Digital Forensics Unit Extraction\nEvidence Hash & Lab Report", C_GOLD),
        ("Stage 6", "Court Trial", "FIR Case File & Challan\nCourt Verdict Tracking (Bail/Conv)", C_GREEN),
    ]

    for idx, (stg, title, desc, col) in enumerate(stages):
        c_left = 0.8 + idx * 1.98
        add_card(s5, c_left, 1.8, 1.85, 4.9, stg, col)
        sb = s5.shapes.add_textbox(Inches(c_left + 0.1), Inches(2.4), Inches(1.65), Inches(4.1))
        stf = sb.text_frame
        stf.word_wrap = True
        
        sp1 = stf.paragraphs[0]
        sp1.text = title
        sp1.font.name = "Segoe UI"
        sp1.font.size = Pt(14)
        sp1.font.bold = True
        sp1.font.color.rgb = C_WHITE
        
        sp2 = stf.add_paragraph()
        sp2.text = desc
        sp2.font.name = "Segoe UI"
        sp2.font.size = Pt(11)
        sp2.font.color.rgb = C_SLATE
        sp2.space_before = Pt(10)

    # ==========================================
    # SLIDE 6: STAGE 1 & 2 DETAILS (REGISTRATION & SCRUTINY)
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s6)
    add_header(s6, "Stage 1 & 2: Registration, Receipt & Scrutiny", "Citizen Experience · Quick Turnaround · Automated Notifications")

    add_card(s6, 0.8, 1.7, 5.75, 5.25, "📝 Stage 1: Front Desk Registration", C_CYAN)
    s1_box = s6.shapes.add_textbox(Inches(1.05), Inches(2.35), Inches(5.25), Inches(4.3))
    s1_tf = s1_box.text_frame
    s1_tf.word_wrap = True
    s1_items = [
        "Walk-in / Online Filing: Operator enters complaint details or drops verification PDF report.",
        "80mm Thermal Complaint Slip: Prints official 80mm thermal receipt with circular NCCIA emblem, Complaint No, Complainant CNIC, and assigned Verification Officer.",
        "Citizen WhatsApp Integration: Instantly transmits WhatsApp receipt with tracking URL directly to citizen's mobile.",
        "Live Tracking QR Code: Citizen can scan the QR code to check status on official portal."
    ]
    for idx, item in enumerate(s1_items):
        p = s1_tf.paragraphs[0] if idx == 0 else s1_tf.add_paragraph()
        p.text = "✔ " + item
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(10)

    add_card(s6, 6.78, 1.7, 5.75, 5.25, "⚖️ Stage 2: Circle Incharge Scrutiny", C_GOLD)
    s2_box = s6.shapes.add_textbox(Inches(7.03), Inches(2.35), Inches(5.25), Inches(4.3))
    s2_tf = s2_box.text_frame
    s2_tf.word_wrap = True
    s2_items = [
        "Jurisdiction Verification: Assistant Director / DSP checks PECA 2016 / PPC crime sections.",
        "Direct VO Assignment: Assigns complaint to field Verification Officer (VO) with priority tagging (High / Medium / Low).",
        "Scrutiny Log & Audit: Tracks time-to-assignment and scrutiny remarks in the permanent case audit trail."
    ]
    for idx, item in enumerate(s2_items):
        p = s2_tf.paragraphs[0] if idx == 0 else s2_tf.add_paragraph()
        p.text = "✔ " + item
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(10)

    # ==========================================
    # SLIDE 7: STAGE 3 & 4 (VERIFICATION & ENQUIRY)
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s7)
    add_header(s7, "Stage 3 & 4: Field Verification & Formal Enquiry", "Evidence-Based Findings · Legal Notices · Case Diary Management")

    add_card(s7, 0.8, 1.7, 5.75, 5.25, "🔍 Stage 3: Verification & 4 Outcomes", C_GREEN)
    v_box = s7.shapes.add_textbox(Inches(1.05), Inches(2.35), Inches(5.25), Inches(4.3))
    vtf = v_box.text_frame
    vtf.word_wrap = True
    v_items = [
        "Preliminary Field Inquiry: VO calls complainant, verifies accused bank accounts, and examines call records.",
        "Recommendation 1 — Enquiry Registered: Sufficient digital evidence exists for formal criminal enquiry.",
        "Recommendation 2 — Case Closure: Closed on lack of evidence, non-pursuance, or mutual settlement.",
        "Recommendation 3 — Case Merge: Scammer already under investigation in an existing primary complaint.",
        "Recommendation 4 — Transfer: Matter belongs to another Circle or external agency (FIA/NAB/Police)."
    ]
    for idx, item in enumerate(v_items):
        p = vtf.paragraphs[0] if idx == 0 else vtf.add_paragraph()
        p.text = "🔹 " + item
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(8)

    add_card(s7, 6.78, 1.7, 5.75, 5.25, "📂 Stage 4: Formal Enquiry & Legal Notices", C_CYAN)
    e_box = s7.shapes.add_textbox(Inches(7.03), Inches(2.35), Inches(5.25), Inches(4.3))
    etf = e_box.text_frame
    etf.word_wrap = True
    e_items = [
        "Enquiry Number Allocation: Systematic numbering (e.g. #LHR-E-80519/26).",
        "Statutory Appearance Notices: Generates Section 160 CrPC legal notices in 80mm thermal and A4 format with QR code verification.",
        "Digital Case Diary (Zimni): IO records daily investigation developments, bank responses, and witness statements.",
        "Final CFR Report: Case Final Report submitted to Incharge for FIR registration approval."
    ]
    for idx, item in enumerate(e_items):
        p = etf.paragraphs[0] if idx == 0 else etf.add_paragraph()
        p.text = "🔹 " + item
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(8)

    # ==========================================
    # SLIDE 8: STAGE 5 & 6 (FORENSICS & COURT TRIAL)
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s8)
    add_header(s8, "Stage 5 & 6: Digital Forensics & Court Trial", "Chain of Custody · Lab Extraction · Prosecution & Court Verdicts")

    add_card(s8, 0.8, 1.7, 5.75, 5.25, "💻 Stage 5: Digital Forensic Lab Portal", C_CYAN)
    f_box = s8.shapes.add_textbox(Inches(1.05), Inches(2.35), Inches(5.25), Inches(4.3))
    ftf = f_box.text_frame
    ftf.word_wrap = True
    f_items = [
        "Device Ingestion & Hashing: Physical intake of phones, hard drives, SIM cards with MD5/SHA256 hash calculation.",
        "Isolated Lab Portal: Forensic examiners operate in a segregated portal for tamper-proof data extraction.",
        "Forensic Examination Report: Digital evidence report directly attached to the IO's case file for prosecution."
    ]
    for idx, item in enumerate(f_items):
        p = ftf.paragraphs[0] if idx == 0 else ftf.add_paragraph()
        p.text = "✔ " + item
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(12)

    add_card(s8, 6.78, 1.7, 5.75, 5.25, "⚖️ Stage 6: Case File (FIR) & Court Trial", C_GOLD)
    c_box = s8.shapes.add_textbox(Inches(7.03), Inches(2.35), Inches(5.25), Inches(4.3))
    ctf = c_box.text_frame
    ctf.word_wrap = True
    c_items = [
        "FIR Registration: Formal Case File created upon established cyber crime offense.",
        "Challan Submission: Prosecution branch tracks court hearings, IO appearance dates, and witness summons.",
        "Court Verdicts Recording: Comprehensive tracking of final verdicts (Conviction, Acquittal, Bail Granted/Rejected, Fines)."
    ]
    for idx, item in enumerate(c_items):
        p = ctf.paragraphs[0] if idx == 0 else ctf.add_paragraph()
        p.text = "✔ " + item
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(12)

    # ==========================================
    # SLIDE 9: ENTERPRISE SECURITY & IP INTELLIGENCE
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s9)
    add_header(s9, "Enterprise Security & IP Audit Intelligence", "Anti-Brute Force Protection · Real IP Detection · Zero Exposure")

    sec_cards = [
        ("🚫 Brute-Force Immunity", "6-attempt/min rate limiting on login routes. Automated IP temporary lockouts against password guessing bots.", C_CYAN),
        ("📡 Real IP & Proxy Detection", "Extracts origin ISP IP, Cloudflare `CF-Connecting-IP`, `X-Real-IP`, and detects VPN/Proxy spoofing.", C_GOLD),
        ("🛡️ RCE & Upload Lockdown", "Uploads folder strictly blocks script execution (`.php`, `.sh`, `.exe`). Magic-byte MIME verification.", C_GREEN),
        ("🔒 Zero-Exposure Hardening", "Strict `.htaccess` rules block access to `.env`, `.git`, database dumps, logs, and sensitive config files.", C_CYAN)
    ]
    for idx, (title, desc, col) in enumerate(sec_cards):
        row = idx // 2
        col_idx = idx % 2
        c_left = 0.8 + col_idx * 5.98
        c_top = 1.7 + row * 2.7
        add_card(s9, c_left, c_top, 5.75, 2.45, title, col)
        sb = s9.shapes.add_textbox(Inches(c_left + 0.2), Inches(c_top + 0.65), Inches(5.35), Inches(1.65))
        stf = sb.text_frame
        stf.word_wrap = True
        p = stf.paragraphs[0]
        p.text = desc
        p.font.name = "Segoe UI"
        p.font.size = Pt(12.5)
        p.font.color.rgb = C_WHITE

    # ==========================================
    # SLIDE 10: USER ROLES & ACCESS HIERARCHY
    # ==========================================
    s10 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s10)
    add_header(s10, "Role-Based Access Control (RBAC) Hierarchy", "Strict Data Isolation Across Circles & Nationwide Command")

    roles = [
        ("Desk Operator", "Front Desk Reception", "New Complaints entry, PDF OCR auto-fill, 80mm slip printing, WhatsApp alert dispatch."),
        ("Verification Officer (VO)", "Field Verification Unit", "Conducts preliminary verification, checks call/bank records, submits VO report."),
        ("Enquiry Officer (IO)", "Core Investigation", "Manages formal enquiries, issues Section 160 notices, updates case diaries (Zimni)."),
        ("Circle Incharge (AD/DSP)", "Circle Command", "Assigns cases, approves enquiry recommendations, reviews Circle KPIs."),
        ("Forensic Specialist", "Digital Forensics Lab", "Performs evidence acquisition, hash validation, and forensic report generation."),
        ("Director General / Admin", "National Headquarters", "Nationwide live analytics, cross-circle transfers, audit logs, IP intelligence.")
    ]
    for idx, (r_name, r_sub, r_desc) in enumerate(roles):
        row = idx // 3
        col_idx = idx % 3
        c_left = 0.8 + col_idx * 3.98
        c_top = 1.7 + row * 2.7
        add_card(s10, c_left, c_top, 3.75, 2.45, r_name, C_GOLD if idx == 5 else C_CYAN)
        sb = s10.shapes.add_textbox(Inches(c_left + 0.2), Inches(c_top + 0.6), Inches(3.35), Inches(1.7))
        stf = sb.text_frame
        stf.word_wrap = True
        
        p1 = stf.paragraphs[0]
        p1.text = r_sub
        p1.font.name = "Segoe UI"
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = C_GREEN
        
        p2 = stf.add_paragraph()
        p2.text = r_desc
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = C_WHITE
        p2.space_before = Pt(4)

    # ==========================================
    # SLIDE 11: 100% ZERO-SCROLL RESPONSIVE UX
    # ==========================================
    s11 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s11)
    add_header(s11, "Responsive UX & Zero-Scroll Interface", "Designed for High-Efficiency Operations on Mobile, Tablet & Desktop")

    add_card(s11, 0.8, 1.7, 5.75, 5.25, "📱 Smart High-Density Column System", C_CYAN)
    ux_box = s11.shapes.add_textbox(Inches(1.05), Inches(2.35), Inches(5.25), Inches(4.3))
    uxtf = ux_box.text_frame
    uxtf.word_wrap = True
    ux_items = [
        "7 Streamlined Smart Columns: Consolidated 10 columns into 7 high-density information units (Tracking/Date, Complainant/CNIC, Crime, Status, Progress, Enquiry, Actions).",
        "Zero Horizontal Scrollbar: Table fits 100% within laptop (1366x768) and desktop (1920x1080) screens without clipping.",
        "Clean Monospace CNIC Formatting: Exact `34100-1115454-6` alignment preventing wrapping bugs.",
        "Touch-Friendly Action Targets: 34px-40px button heights with high-contrast color badges for rapid one-click actions."
    ]
    for idx, item in enumerate(ux_items):
        p = uxtf.paragraphs[0] if idx == 0 else uxtf.add_paragraph()
        p.text = "✔ " + item
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(10)

    add_card(s11, 6.78, 1.7, 5.75, 5.25, "🖨️ Multi-Format Thermal & Legal Prints", C_GOLD)
    pr_box = s11.shapes.add_textbox(Inches(7.03), Inches(2.35), Inches(5.25), Inches(4.3))
    prtf = pr_box.text_frame
    prtf.word_wrap = True
    pr_items = [
        "80mm Thermal Receipt Slip: High-speed point-of-sale thermal slip with official NCCIA seal, QR code, and Verification Officer name.",
        "80mm Legal Appearance Notice: Section 160 CrPC witness/accused call notice with date, venue, and reporting officer.",
        "A4 Full Investigation Report: Comprehensive case dossier with Complainant, Accused, Financial Trails, and Attachments.",
        "Unified Visual Identity: Authentic circular NCCIA seal standardized across all prints."
    ]
    for idx, item in enumerate(pr_items):
        p = prtf.paragraphs[0] if idx == 0 else prtf.add_paragraph()
        p.text = "✔ " + item
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(10)

    # ==========================================
    # SLIDE 12: LIVE DEMONSTRATION GUIDE
    # ==========================================
    s12 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s12)
    add_header(s12, "Live System Demonstration Guide", "Step-by-Step Demonstration Flow for Leadership & Stakeholders")

    demo_steps = [
        ("Step 1: AI PDF Auto-Fill", "Open `/complaints/import-pdf` ➔ Drop CCW verification PDF ➔ Watch Name, CNIC, Phone, Offence, and Accused Bank details auto-fill in 1 second."),
        ("Step 2: Instant Slip & WhatsApp", "Click `Save & Register` ➔ Print 80mm thermal slip with VO name & QR code ➔ Trigger automated WhatsApp citizen alert."),
        ("Step 3: Scrutiny & Assignment", "Login as Circle Incharge ➔ Scrutinize complaint ➔ Assign to Verification Officer (VO) with priority tagging."),
        ("Step 4: Field Verification", "Login as VO ➔ Review digital evidence ➔ Recommend Enquiry Registration ➔ Submit report for approval."),
        ("Step 5: Formal Enquiry & Notices", "Login as IO ➔ View registered enquiry `#LHR-E-80519/26` ➔ Generate 80mm/A4 Legal Appearance Notice ➔ Add Case Diary (Zimni)."),
        ("Step 6: Security & IP Intelligence", "Open `/login-history` ➔ Search client IP / Real IP ➔ Demonstrate automated spoofing detection and audit logging.")
    ]
    for idx, (title, desc) in enumerate(demo_steps):
        row = idx // 2
        col_idx = idx % 2
        c_left = 0.8 + col_idx * 5.98
        c_top = 1.7 + row * 1.8
        add_card(s12, c_left, c_top, 5.75, 1.6, title, C_GOLD if idx == 0 else C_CYAN)
        sb = s12.shapes.add_textbox(Inches(c_left + 0.2), Inches(c_top + 0.5), Inches(5.35), Inches(1.0))
        stf = sb.text_frame
        stf.word_wrap = True
        p = stf.paragraphs[0]
        p.text = desc
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_WHITE

    # ==========================================
    # SLIDE 13: CONCLUSION & FUTURE ROADMAP
    # ==========================================
    s13 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s13)
    add_header(s13, "Conclusion & Strategic Impact", "Empowering the Nation with Cutting-Edge Cyber Defense Technology")

    add_card(s13, 0.8, 1.7, 11.733, 5.25, "🌟 Transformational Benefits for NCCIA", C_GOLD)
    cb_box = s13.shapes.add_textbox(Inches(1.1), Inches(2.4), Inches(11.133), Inches(4.3))
    cbtf = cb_box.text_frame
    cbtf.word_wrap = True
    c_points = [
        "10x Faster Case Processing: Registration time reduced from 20 minutes to under 60 seconds with AI OCR auto-fill.",
        "Zero Paper Loss: Complete digital lifecycle ensures evidence and case diaries are permanently preserved and searchable.",
        "100% Institutional Transparency: Complete audit trail with timestamps, IP intelligence, and officer accountability.",
        "Nationwide High-Availability: Optimized architecture capable of handling 10 Million+ concurrent requests and 1 Crore records seamlessly.",
        "Citizen Trust & Satisfaction: Instant WhatsApp receipts, QR verification, and automated status updates empower public confidence."
    ]
    for idx, pt in enumerate(c_points):
        p = cbtf.paragraphs[0] if idx == 0 else cbtf.add_paragraph()
        p.text = "⭐ " + pt
        p.font.name = "Segoe UI"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(14)

    # Save presentation
    output_path = os.path.abspath("f:/NCCIA/NCCIA_Executive_Master_Presentation.pptx")
    prs.save(output_path)
    print(f"Master presentation successfully generated at: {output_path}")

if __name__ == "__main__":
    create_master_presentation()
